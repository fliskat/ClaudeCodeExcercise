"""
Adds a Future State Architecture slide to Presentation5.pptx.
Layout: Source Systems → Incorta Connectors → Incorta (central) → Consumption & AI
ACDC Schema includes two new boxes: Golden Record & Customer Hierarchy
SQL and Python Script labels on connector arrows.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy

INPUT  = "/Users/ratnadeep.bose/Downloads/Epredia Temp Doc/Presentation5.pptx"
OUTPUT = "/Users/ratnadeep.bose/Downloads/Epredia Temp Doc/Epredia_Future_State_Architecture.pptx"

# ── Colour palette (Midnight Executive) ──────────────────────────────────────
C_NAVY      = RGBColor(0x1E, 0x27, 0x61)   # dominant navy
C_TEAL      = RGBColor(0x02, 0x80, 0x90)   # Incorta accent
C_ICE       = RGBColor(0xCA, 0xDC, 0xFC)   # light fill
C_GOLD      = RGBColor(0xF5, 0xA6, 0x23)   # SQL/script connector labels
C_GREEN     = RGBColor(0x1A, 0x7A, 0x4A)   # new boxes (Golden Record / Cust Hier)
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK_GREY = RGBColor(0x2B, 0x2B, 0x2B)
C_LIGHT_BG  = RGBColor(0xF0, 0xF4, 0xFF)
C_SEC_BAR   = RGBColor(0x36, 0x45, 0x4F)


def rgb(color: RGBColor):
    return f"{color.rgb:06X}"


def add_rect(slide, l, t, w, h, fill_color, border_color=None, border_pt=1.5, radius=None):
    """Add a rectangle shape and return it."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    if radius:
        # Add rounded corners via XML
        sp = shape.element
        spPr = sp.find(qn("p:spPr"))
        prstGeom = spPr.find(qn("a:prstGeom"))
        if prstGeom is not None:
            prstGeom.attrib["prst"] = "roundRect"
            avLst = prstGeom.find(qn("a:avLst"))
            if avLst is None:
                avLst = etree.SubElement(prstGeom, qn("a:avLst"))
            gd = etree.SubElement(avLst, qn("a:gd"))
            gd.set("name", "adj")
            gd.set("fmla", f"val {radius}")
    return shape


def set_text(shape, text, font_size=9, bold=False, color=C_WHITE, align=PP_ALIGN.CENTER, wrap=True):
    tf = shape.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    from pptx.util import Pt
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    # Vertical centre
    from pptx.enum.text import MSO_ANCHOR
    tf.auto_size = None
    tf.word_wrap = wrap
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_label(slide, l, t, w, h, text, font_size=7.5, bold=False,
               color=C_DARK_GREY, align=PP_ALIGN.CENTER):
    """Transparent label box."""
    shape = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=C_NAVY, width_pt=1.5):
    """Draw a connector line with arrowhead."""
    from pptx.util import Emu
    cx = slide.shapes._spTree

    # Build freeform connector in XML
    x1e, y1e = int(Inches(x1)), int(Inches(y1))
    x2e, y2e = int(Inches(x2)), int(Inches(y2))
    w = abs(x2e - x1e) or 1
    h = abs(y2e - y1e) or 1
    off_x = min(x1e, x2e)
    off_y = min(y1e, y2e)

    # local coords
    lx1 = x1e - off_x
    ly1 = y1e - off_y
    lx2 = x2e - off_x
    ly2 = y2e - off_y

    clr = str(color)
    w_emu = int(Pt(width_pt).pt * 12700)

    xml = f"""<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
              xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
              xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr>
    <p:cNvPr id="9999" name="arrow"/>
    <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm>
      <a:off x="{off_x}" y="{off_y}"/>
      <a:ext cx="{w}" cy="{h}"/>
    </a:xfrm>
    <a:prstGeom prst="line"><a:avLst/></a:prstGeom>
    <a:ln w="{w_emu}">
      <a:solidFill><a:srgbClr val="{clr}"/></a:solidFill>
      <a:tailEnd type="arrow"/>
    </a:ln>
  </p:spPr>
  <p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>
</p:sp>"""
    sp = etree.fromstring(xml)
    cx.append(sp)


# ── Build slide ───────────────────────────────────────────────────────────────
prs = Presentation(INPUT)
slide_layout = prs.slide_layouts[6]  # blank
slide = prs.slides.add_slide(slide_layout)

W, H = 13.33, 7.5   # inches (16:9 widescreen)

# Background
bg = add_rect(slide, 0, 0, W, H, C_LIGHT_BG)

# ── Title bar ──────────────────────────────────────────────────────────────
title_bar = add_rect(slide, 0, 0, W, 0.55, C_NAVY)
add_label(slide, 0.15, 0.08, 11, 0.42,
          "Epredia Future State Data Architecture — Incorta as Central Platform",
          font_size=13, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)

# ── Column headers ────────────────────────────────────────────────────────
col_headers = [
    (0.15,  "SOURCE SYSTEMS"),
    (2.85,  "INCORTA CONNECTORS"),
    (5.25,  "INCORTA  (Central Platform)"),
    (9.75,  "CONSUMPTION & AI"),
]
for x, label in col_headers:
    hdr = add_rect(slide, x, 0.65, 2.5 if x < 5 else (4.3 if x == 5.25 else 3.4), 0.32, C_NAVY, radius=20000)
    set_text(hdr, label, font_size=7.5, bold=True, color=C_WHITE)

# Fix widths per column
add_rect(slide, 0.15,  0.65, 2.55, 0.32, C_NAVY, radius=20000)
add_rect(slide, 2.85,  0.65, 2.25, 0.32, C_NAVY, radius=20000)
add_rect(slide, 5.25,  0.65, 4.30, 0.32, C_NAVY, radius=20000)
add_rect(slide, 9.75,  0.65, 3.40, 0.32, C_NAVY, radius=20000)

# ── SOURCE SYSTEMS (col 1) ────────────────────────────────────────────────
sources = [
    ("ERP\n(SAP)",                     C_NAVY),
    ("CRM\n(Salesforce / ServiceMax)", C_NAVY),
    ("Internal Data\n(SharePoint)",    C_NAVY),
    ("Files & APIs\n(External Data)",  C_NAVY),
    ("Email\nAttachments",             C_NAVY),
]
src_x, src_w, src_h = 0.18, 2.45, 0.70
src_y_start = 1.10
src_gap = 0.85

for i, (label, color) in enumerate(sources):
    y = src_y_start + i * src_gap
    b = add_rect(slide, src_x, y, src_w, src_h, color, C_ICE, border_pt=1.2, radius=30000)
    set_text(b, label, font_size=9, bold=True, color=C_WHITE)

# ── INCORTA CONNECTORS (col 2) ────────────────────────────────────────────
connectors = [
    "SAP Connector\n(Marketplace)",
    "Salesforce Connector\n(Marketplace)",
    "SharePoint Connector\n(Marketplace)",
    "REST / API Connector\n(Python Script)",
    "Power Automate\n→ SharePoint → Incorta",
]
con_x, con_w, con_h = 2.88, 2.20, 0.70

for i, label in enumerate(connectors):
    y = src_y_start + i * src_gap
    b = add_rect(slide, con_x, y, con_w, con_h, C_TEAL, C_WHITE, border_pt=1.0, radius=25000)
    set_text(b, label, font_size=8, bold=False, color=C_WHITE)

# ── INCORTA CENTRAL BOX ───────────────────────────────────────────────────
inc_x, inc_y, inc_w, inc_h = 5.28, 1.00, 4.28, 5.90
incorta_bg = add_rect(slide, inc_x, inc_y, inc_w, inc_h, C_WHITE, C_TEAL, border_pt=2.5, radius=18000)

# Incorta label at top of central box
inc_title = add_rect(slide, inc_x + 0.05, inc_y + 0.05, inc_w - 0.10, 0.40, C_TEAL, radius=15000)
set_text(inc_title, "INCORTA", font_size=13, bold=True, color=C_WHITE)

# ── ACDC Schema sub-box ───────────────────────────────────────────────────
acdc_x, acdc_y = inc_x + 0.15, inc_y + 0.58
acdc_w, acdc_h = inc_w - 0.30, 2.65
acdc_box = add_rect(slide, acdc_x, acdc_y, acdc_w, acdc_h, C_ICE, C_NAVY, border_pt=1.5, radius=12000)
add_label(slide, acdc_x + 0.08, acdc_y + 0.04, acdc_w - 0.16, 0.25,
          "ACDC Schema", font_size=8.5, bold=True, color=C_NAVY)

# Raw Data Tables box
rdt = add_rect(slide, acdc_x + 0.12, acdc_y + 0.32, acdc_w - 0.24, 0.46,
               C_NAVY, C_WHITE, border_pt=1.0, radius=10000)
set_text(rdt, "Raw Data Tables", font_size=9, bold=True, color=C_WHITE)

# Golden Record (NEW - smaller)
gr = add_rect(slide, acdc_x + 0.12, acdc_y + 0.88, (acdc_w - 0.36) / 2, 0.55,
              C_GREEN, C_WHITE, border_pt=1.2, radius=10000)
set_text(gr, "Golden\nRecord", font_size=8.5, bold=True, color=C_WHITE)

# Customer Hierarchy (NEW - smaller)
ch = add_rect(slide, acdc_x + 0.12 + (acdc_w - 0.36) / 2 + 0.12,
              acdc_y + 0.88, (acdc_w - 0.36) / 2, 0.55,
              C_GREEN, C_WHITE, border_pt=1.2, radius=10000)
set_text(ch, "Customer\nHierarchy", font_size=8.5, bold=True, color=C_WHITE)

# NEW badges
add_label(slide, acdc_x + 0.12,                            acdc_y + 0.88 - 0.18, 1.2, 0.20,
          "★ NEW", font_size=7, bold=True, color=C_GREEN)
add_label(slide, acdc_x + 0.12 + (acdc_w - 0.36) / 2 + 0.12, acdc_y + 0.88 - 0.18, 1.2, 0.20,
          "★ NEW", font_size=7, bold=True, color=C_GREEN)

# Label: Python Script under golden record boxes
add_label(slide, acdc_x + 0.08, acdc_y + 1.50, acdc_w - 0.16, 0.22,
          "populated via Python Script / PySpark", font_size=7, bold=False, color=C_TEAL,
          align=PP_ALIGN.CENTER)

# ── Business Schema box ───────────────────────────────────────────────────
bs_y = inc_y + 0.58 + acdc_h + 0.12
bs = add_rect(slide, inc_x + 0.15, bs_y, inc_w - 0.30, 0.50, C_NAVY, C_WHITE, border_pt=1.0, radius=10000)
set_text(bs, "Business Schema", font_size=9, bold=True, color=C_WHITE)
add_label(slide, inc_x + 0.15, bs_y + 0.52, inc_w - 0.30, 0.18,
          "SQL Views & Transformations", font_size=7, bold=False, color=C_TEAL, align=PP_ALIGN.CENTER)

# ── Materialized Views box ────────────────────────────────────────────────
mv_y = bs_y + 0.75
mv = add_rect(slide, inc_x + 0.15, mv_y, inc_w - 0.30, 0.50, RGBColor(0x02, 0x50, 0x70), C_WHITE, border_pt=1.0, radius=10000)
set_text(mv, "Materialized Views  (PySpark / SQL)", font_size=9, bold=True, color=C_WHITE)

# ── Presentation Tier box ─────────────────────────────────────────────────
pt_y = mv_y + 0.65
pt = add_rect(slide, inc_x + 0.15, pt_y, inc_w - 0.30, 0.45, C_TEAL, C_WHITE, border_pt=1.0, radius=10000)
set_text(pt, "Incorta Presentation Tier", font_size=9, bold=True, color=C_WHITE)

# ── CONSUMPTION & AI (col 4) ──────────────────────────────────────────────
consumers = [
    ("Dashboards &\nSelf-Service BI",  C_NAVY),
    ("Power BI\nReports & Dashboard",  C_NAVY),
    ("AgentForce\n(Salesforce AI)",    C_NAVY),
    ("Salesforce\nAPI Integration",    C_NAVY),
]
cons_x, cons_w, cons_h = 9.78, 3.35, 0.72
cons_y_start = 1.10

for i, (label, color) in enumerate(consumers):
    y = cons_y_start + i * src_gap
    b = add_rect(slide, cons_x, y, cons_w, cons_h, color, C_ICE, border_pt=1.2, radius=30000)
    set_text(b, label, font_size=9, bold=True, color=C_WHITE)

# ── SECURITY bar at bottom ────────────────────────────────────────────────
sec = add_rect(slide, 0.15, 7.02, W - 0.30, 0.32, C_SEC_BAR)
set_text(sec, "🔒  Security · Governance · Data Lineage · Metadata Management",
         font_size=8.5, bold=False, color=C_WHITE)

# ── CONNECTOR ARROWS & SQL / PYTHON SCRIPT LABELS ────────────────────────
# Source → Connector arrows (horizontal)
for i in range(5):
    y_mid = src_y_start + i * src_gap + src_h / 2
    add_arrow(slide, src_x + src_w, y_mid, con_x, y_mid, color=C_NAVY, width_pt=1.2)

# Connector → Incorta central arrows + SQL label
for i in range(5):
    y_mid = src_y_start + i * src_gap + src_h / 2
    add_arrow(slide, con_x + con_w, y_mid, inc_x, y_mid, color=C_TEAL, width_pt=1.5)
    # SQL label on the arrow midpoint
    lbl_x = con_x + con_w + 0.03
    add_label(slide, lbl_x, y_mid - 0.25, 0.75, 0.22,
              "SQL", font_size=6.5, bold=True, color=C_GOLD)

# Incorta → Consumption arrows + Python Script label
cons_centers = [cons_y_start + i * src_gap + cons_h / 2 for i in range(4)]
inc_right_x = inc_x + inc_w
for i, y_mid in enumerate(cons_centers):
    add_arrow(slide, inc_right_x, y_mid, cons_x, y_mid, color=C_NAVY, width_pt=1.5)
    add_label(slide, inc_right_x + 0.03, y_mid - 0.25, 1.05, 0.22,
              "Python / SQL", font_size=6.5, bold=True, color=C_GOLD)

# ── Internal arrows within Incorta (ACDC → Business Schema → MV → Pres Tier)
mid_x = inc_x + inc_w / 2
# ACDC bottom → Business Schema
acdc_bottom = acdc_y + acdc_h
add_arrow(slide, mid_x, acdc_bottom, mid_x, bs_y, color=C_TEAL, width_pt=1.2)
# Business Schema → MV
add_arrow(slide, mid_x, bs_y + 0.50, mid_x, mv_y, color=C_TEAL, width_pt=1.2)
# MV → Presentation Tier
add_arrow(slide, mid_x, mv_y + 0.50, mid_x, pt_y, color=C_TEAL, width_pt=1.2)

# ── Save ──────────────────────────────────────────────────────────────────
prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
